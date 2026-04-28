"""
Generate MINDSET Bolivia Assessment Presentation
Focused on step-by-step drought/flood analysis workflow and user inputs.
Does NOT reference SiSePuede (SSP) files — those are a separate adaptation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x95, 0x95, 0x95)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
VERY_LIGHT_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
LIGHT_GREEN_BG = RGBColor(0xE8, 0xF6, 0xEF)
LIGHT_ORANGE_BG = RGBColor(0xFD, 0xF2, 0xE9)
LIGHT_RED_BG = RGBColor(0xFD, 0xED, 0xED)
NEUTRAL_BG = RGBColor(0xE8, 0xEA, 0xED)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_frame(slide, left, top, width, height, items, font_size=13,
                     color=DARK_GRAY, spacing=Pt(6), bold_prefix=False,
                     bullet_char="\u2022 "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        prefixed = bullet_char + item if bullet_char else item

        if bold_prefix and ': ' in prefixed:
            colon_pos = prefixed.index(': ')
            part1 = prefixed[:colon_pos + 2]
            part2 = prefixed[colon_pos + 2:]
            run1 = p.add_run()
            run1.text = part1
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = 'Calibri'
            run2 = p.add_run()
            run2.text = part2
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
            run2.font.name = 'Calibri'
        else:
            run = p.add_run()
            run.text = prefixed
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = 'Calibri'
    return tf


def add_footer(slide, text="MINDSET Assessment for Bolivia"):
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(5), Inches(0.4),
                 text, font_size=9, color=MED_GRAY)


def add_source_line(slide, text):
    add_text_box(slide, Inches(0.7), Inches(6.85), Inches(11.5), Inches(0.4),
                 text, font_size=9, color=MED_GRAY)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(4.5), DARK_BLUE)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             "MINDSET Model", font_size=42, bold=True, color=WHITE)
add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(1),
             "Capability Assessment for Bolivia",
             font_size=32, bold=False, color=RGBColor(0xA0, 0xCF, 0xED))
add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
             "ENSO Climate Shocks & Energy Sector Vulnerabilities",
             font_size=18, color=RGBColor(0x80, 0xBB, 0xDD))

add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             "Two questions from the research team:",
             font_size=16, bold=True, color=DARK_BLUE)

add_bullet_frame(slide, Inches(1.3), Inches(5.7), Inches(10.5), Inches(1.5), [
    "Can MINDSET help assess ENSO-related droughts and floods \u2014 their macroeconomic and distributional impacts?",
    "Can MINDSET help analyze declining gas reserves and hydropower uncertainty \u2014 energy shocks and low-carbon pathways?"
], font_size=14, color=DARK_GRAY, bullet_char="")

# ============================================================
# SLIDE 2: Bottom Line Up Front
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "Bottom Line", font_size=28, bold=True, color=WHITE)
add_footer(slide)

# Task 1 box
add_shape_bg(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.2), VERY_LIGHT_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.5),
             "Task #1: ENSO / Floods / Droughts",
             font_size=18, bold=True, color=DARK_BLUE)

add_shape_bg(slide, Inches(1.0), Inches(2.3), Inches(3.5), Inches(0.45), ACCENT_GREEN)
add_text_box(slide, Inches(1.0), Inches(2.32), Inches(3.5), Inches(0.4),
             "  WELL SUITED  \u2014  Moderate-High Confidence",
             font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_bullet_frame(slide, Inches(1.0), Inches(3.0), Inches(5.2), Inches(3.5), [
    "Floods and droughts are explicitly cited as motivating use cases in MINDSET documentation (Capital damages module).",
    "Supply constraint channel propagates sector-level damage through the full IO system \u2014 GDP, employment, prices, household welfare.",
    "Employment by skill (low/high) and gender (female/male); household consumption by income group.",
    "Bolivia is directly represented in the GLORIA MRIO database (164 countries, 120 sectors, base year 2019).",
    "Key requirement: users must provide sector-level damage estimates. MINDSET models the economic propagation, not the physical hazard."
], font_size=12, color=DARK_GRAY, spacing=Pt(8))

# Task 2 box
add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), LIGHT_ORANGE_BG)
add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.5),
             "Task #2: Gas Decline & Hydropower Uncertainty",
             font_size=18, bold=True, color=DARK_BLUE)

add_shape_bg(slide, Inches(7.1), Inches(2.3), Inches(3.5), Inches(0.45), ACCENT_ORANGE)
add_text_box(slide, Inches(7.1), Inches(2.32), Inches(3.5), Inches(0.4),
             "  PARTIAL FIT  \u2014  Moderate-Low Confidence",
             font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_bullet_frame(slide, Inches(7.1), Inches(3.0), Inches(5.2), Inches(3.5), [
    "Can model supply constraints on gas extraction (sector 27) and electricity (sector 63), plus interfuel substitution via price elasticities.",
    "Useful for macro impact of externally defined energy shocks (GDP, employment, supply chain effects, household welfare).",
    "Cannot model resource depletion trajectories, hydropower variability, energy system optimization, or long-run low-carbon investment pathways.",
    "Best role: complementary macro/distributional analysis alongside a dedicated energy system model (TIMES, OSeMOSYS).",
    "Electricity sector is a single aggregate \u2014 no hydro vs. thermal vs. solar distinction."
], font_size=12, color=DARK_GRAY, spacing=Pt(8))

# ============================================================
# SLIDE 3: What is MINDSET (brief)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "What is MINDSET?", font_size=28, bold=True, color=WHITE)
add_footer(slide)

add_text_box(slide, Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.6),
             "A Multi-Regional Input-Output (MRIO) model built on the GLORIA database. "
             "It traces how a shock in one sector or country propagates through global supply chains "
             "to affect output, employment, prices, and household welfare across all sectors and trading partners.",
             font_size=14, color=DARK_GRAY)

# Three columns
col_data = [
    ("Core Design", VERY_LIGHT_BLUE, [
        "Demand-side: Leontief inverse L = (I\u2212A)\u207B\u00B9",
        "Supply-side: Ghosh inverse G = (I\u2212B)\u207B\u00B9 with criticality weighting",
        "Not a CGE \u2014 allows quantity rationing, no market clearing",
        "Iterative convergence loop on household income and tax revenue",
        "Short-to-medium-term shock assessment",
    ]),
    ("Data Coverage", LIGHT_GREEN_BG, [
        "164 countries (GLORIA v57, base year 2019)",
        "120 production sectors (agriculture through services)",
        "Full bilateral trade flows between all countries",
        "Employment data by skill level and gender",
        "Household elasticities by income group (USDA)",
    ]),
    ("Disaster-Relevant Channels", LIGHT_ORANGE_BG, [
        "Supply constraints (the primary disaster shock input)",
        "Exogenous investment (reconstruction spending)",
        "Exogenous final demand (post-disaster demand changes)",
        "Price changes (scarcity-driven price spikes)",
        "Technology substitution (interfuel switching)",
    ]),
]

for i, (title, bg, items) in enumerate(col_data):
    x = Inches(0.7) + i * Inches(4.1)
    w = Inches(3.9)
    add_shape_bg(slide, x, Inches(2.3), w, Inches(4.2), bg)
    add_text_box(slide, x + Inches(0.2), Inches(2.4), w - Inches(0.4), Inches(0.4),
                 title, font_size=15, bold=True, color=DARK_BLUE)
    add_bullet_frame(slide, x + Inches(0.2), Inches(2.9), w - Inches(0.4), Inches(3.4),
                     items, font_size=11, color=DARK_GRAY, spacing=Pt(5))

add_source_line(slide,
    "Sources: README.md; Documentation/Capital damages.md; "
    "Documentation/Supply constraint.md; SourceCode/InputOutput.py")

# ============================================================
# SLIDE 4: Step-by-Step Pipeline Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "Estimating Drought/Flood Impacts: Step-by-Step Pipeline",
             font_size=26, bold=True, color=WHITE)
add_footer(slide)

# Five pipeline steps as boxes with arrows
steps = [
    ("Step 1",
     "Characterize the\nPhysical Hazard",
     "Identify which sectors\nand regions are affected\nby the flood/drought.\nQuantify capital destruction\nor output loss (%).",
     LIGHT_RED_BG,
     "USER\nPROVIDES",
     ACCENT_RED),
    ("Step 2",
     "Translate Damage\nto Supply Constraints",
     "Map physical damage to\nGLORIA sector codes (1\u2013120).\nExpress as % output reduction.\nHallegatte (2014): \u0394Y \u2248 \u0394K.",
     LIGHT_ORANGE_BG,
     "USER\nPROVIDES",
     ACCENT_ORANGE),
    ("Step 3",
     "Build Scenario\nExcel Workbook",
     "Fill the \"Supply constraint\"\nsheet with: Country ISO,\nProduct code, Value, Type.\nOptional: Final demand,\nInvestment sheets.",
     VERY_LIGHT_BLUE,
     "USER\nPROVIDES",
     LIGHT_BLUE),
    ("Step 4",
     "Run MINDSET\n(RunMINDSET.py)",
     "Model builds L and G\ninverses, propagates\nshocks through supply\nchains, iterates until\nconvergence.",
     LIGHT_GREEN_BG,
     "MODEL\nRUNS",
     ACCENT_GREEN),
    ("Step 5",
     "Analyze Results\n(8-sheet workbook)",
     "Output, employment, GDP,\ndemand, household welfare,\nprices, tax revenue,\nemissions \u2014 all by sector\nand impact channel.",
     NEUTRAL_BG,
     "MODEL\nOUTPUTS",
     MED_BLUE),
]

box_w = Inches(2.2)
gap = Inches(0.18)
start_x = Inches(0.5)
arrow_w = Inches(0.22)

for i, (step_label, title, desc, bg, badge_text, badge_color) in enumerate(steps):
    x = start_x + i * (box_w + gap + arrow_w)

    # Main box
    add_shape_bg(slide, x, Inches(1.5), box_w, Inches(5.0), bg)

    # Badge at top
    add_shape_bg(slide, x, Inches(1.5), box_w, Inches(0.55), badge_color)
    add_text_box(slide, x, Inches(1.52), box_w, Inches(0.5),
                 badge_text, font_size=9, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Step number
    add_text_box(slide, x + Inches(0.15), Inches(2.15), box_w - Inches(0.3), Inches(0.3),
                 step_label, font_size=11, bold=True, color=MED_GRAY)

    # Title
    add_text_box(slide, x + Inches(0.15), Inches(2.45), box_w - Inches(0.3), Inches(0.8),
                 title, font_size=14, bold=True, color=DARK_BLUE)

    # Description
    add_text_box(slide, x + Inches(0.15), Inches(3.4), box_w - Inches(0.3), Inches(2.8),
                 desc, font_size=11, color=DARK_GRAY)

    # Arrow between boxes
    if i < len(steps) - 1:
        ax = x + box_w + Inches(0.01)
        add_text_box(slide, ax, Inches(3.5), gap + arrow_w, Inches(0.5),
                     "\u25B6", font_size=22, color=LIGHT_BLUE,
                     alignment=PP_ALIGN.CENTER)

add_source_line(slide,
    "Sources: Documentation/Capital damages.md; Documentation/Supply constraint.md; "
    "RunMINDSET.py; SourceCode/scenario.py; SourceCode/InputOutput.py")

# ============================================================
# SLIDE 5: What Users Must Provide (Inputs)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "What Users Must Provide: The Scenario Excel Workbook",
             font_size=26, bold=True, color=WHITE)
add_footer(slide)

# --- Primary input: Supply constraint sheet ---
add_shape_bg(slide, Inches(0.7), Inches(1.3), Inches(7.5), Inches(5.5), VERY_LIGHT_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(7.0), Inches(0.4),
             "Primary input: \"Supply constraint\" sheet  (required for flood/drought)",
             font_size=15, bold=True, color=DARK_BLUE)

add_text_box(slide, Inches(1.0), Inches(1.9), Inches(7.0), Inches(0.4),
             "Data starts at row 15 (14 header rows are skipped). Four columns:",
             font_size=12, color=DARK_GRAY)

# Column table
col_headers = ["Column Name", "Description", "Example"]
col_w_list = [Inches(1.6), Inches(2.8), Inches(2.6)]
col_x_list = [Inches(1.0), Inches(2.6), Inches(5.4)]
y_row = Inches(2.4)
rh = Inches(0.35)

# Header row
for hdr, cx, cw in zip(col_headers, col_x_list, col_w_list):
    add_shape_bg(slide, cx, y_row, cw, rh, DARK_BLUE)
    add_text_box(slide, cx + Inches(0.08), y_row + Inches(0.03), cw - Inches(0.15), rh,
                 hdr, font_size=11, bold=True, color=WHITE)

table_rows = [
    ("Country ISO*", "Country ISO3 code (or \"ALL\")", "BOL"),
    ("Product code*", "GLORIA sector 1\u2013120, range, or list", "1-10  or  5,10,15"),
    ("Value*", "Output reduction as negative decimal", "-0.30  (= 30% loss)"),
    ("Type*", "Currently only \"rel\" (relative)", "rel"),
]

for ri, (c1, c2, c3) in enumerate(table_rows):
    y = y_row + (ri + 1) * rh
    bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
    for text, cx, cw in zip([c1, c2, c3], col_x_list, col_w_list):
        add_shape_bg(slide, cx, y, cw, rh, bg)
        is_code = (text == c1)
        add_text_box(slide, cx + Inches(0.08), y + Inches(0.03), cw - Inches(0.15), rh,
                     text, font_size=10,
                     bold=is_code,
                     color=DARK_BLUE if is_code else DARK_GRAY,
                     font_name='Consolas' if is_code else 'Calibri')

# Example scenario
add_text_box(slide, Inches(1.0), Inches(4.3), Inches(7.0), Inches(0.4),
             "Example: modeling a severe drought in Bolivia",
             font_size=13, bold=True, color=DARK_BLUE)

example_text = (
    "Country ISO*    Product code*    Value*    Type*\n"
    "BOL             1-10             -0.35     rel        \u2190 35% loss in agriculture (sectors 1\u201310)\n"
    "BOL             11-23            -0.15     rel        \u2190 15% loss in remaining agriculture\n"
    "BOL             63               -0.10     rel        \u2190 10% loss in electricity (hydropower)\n"
    "BOL             93-96            -0.05     rel        \u2190 5% loss in water/utilities"
)
add_text_box(slide, Inches(1.0), Inches(4.75), Inches(7.0), Inches(1.8),
             example_text, font_size=10, color=DARK_GRAY, font_name='Consolas')

# --- Optional sheets ---
add_shape_bg(slide, Inches(8.5), Inches(1.3), Inches(4.1), Inches(5.5), LIGHT_GREEN_BG)
add_text_box(slide, Inches(8.7), Inches(1.4), Inches(3.7), Inches(0.4),
             "Optional sheets", font_size=15, bold=True, color=DARK_BLUE)

optional_items = [
    "Final demand: reduce household or government consumption post-disaster. Columns: Producing/Consuming country, Product code, FD code (FD_1=HH, FD_3=Gov, FD_4=GFCF), Value, Type.",
    "Investment by: reconstruction or recovery investment by sector. Columns: Country ISO, Sector investing code, Value (abs-b), Type.",
    "Price change: non-tax cost shocks (scarcity). Columns: Origin/Target country, Origin/Target sector, Value, Type (rel or tax-abs).",
    "IO coefficients: alter production recipes. Columns: Origin/Target country, Origin/Target sector, Value, Type (replace).",
    "Carbon tax: not needed for ENSO analysis.",
    "CBAM: not needed for ENSO analysis.",
]

add_bullet_frame(slide, Inches(8.7), Inches(1.9), Inches(3.7), Inches(4.7),
                 optional_items, font_size=10, color=DARK_GRAY, spacing=Pt(6),
                 bold_prefix=True)

add_source_line(slide,
    "Sources: SourceCode/scenario.py (lines 556\u2013577: set_supply_constraint; "
    "lines 209\u2013313: set_exog_fd; lines 321\u2013383: set_exog_inv)")

# ============================================================
# SLIDE 6: From Physical Damage to Supply Constraint
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "From Physical Damage to Supply Constraint: The Translation Step",
             font_size=26, bold=True, color=WHITE)
add_footer(slide)

# Left: Hallegatte framework
add_shape_bg(slide, Inches(0.7), Inches(1.3), Inches(6.0), Inches(3.2), VERY_LIGHT_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5.5), Inches(0.4),
             "Hallegatte (2014a): Capital Damage \u2192 Output Loss",
             font_size=15, bold=True, color=DARK_BLUE)

add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.5), Inches(2.4),
             "The theoretical relationship (modified Cobb-Douglas):\n\n"
             "    \u0394Y(t\u2080) = (1/\u03bc) \u00b7 r \u00b7 \u0394K\n\n"
             "where \u03bc is capital elasticity, r is interest rate + depreciation, "
             "\u0394K is capital destroyed. With standard parameter values, "
             "output loss is roughly 3\u00d7 the asset value destroyed.\n\n"
             "Simplified assumption (Hallegatte 2008): productive asset loss "
             "to output loss is approximately 1:1. This is the more "
             "conservative and commonly used mapping.",
             font_size=12, color=DARK_GRAY)

# Right: practical guidance
add_shape_bg(slide, Inches(7.0), Inches(1.3), Inches(5.6), Inches(3.2), LIGHT_GREEN_BG)
add_text_box(slide, Inches(7.3), Inches(1.4), Inches(5.0), Inches(0.4),
             "Practical Guidance for Users",
             font_size=15, bold=True, color=DARK_BLUE)

add_bullet_frame(slide, Inches(7.3), Inches(1.9), Inches(5.0), Inches(2.4), [
    "Estimate capital destruction by sector from external hazard models, historical data, or damage assessments (e.g., DesInventar, EM-DAT, post-disaster surveys).",
    "Map damaged sectors to GLORIA sector codes (1\u2013120). Agriculture = sectors 1\u201323; mining = 24\u201329; electricity = 63; water = 93\u201396; construction = 107.",
    "Express as percentage output reduction: 25% capital destruction \u2192 Value = -0.25 in the Supply constraint sheet.",
    "For floods, consider: physical infrastructure damage + crop losses + transport disruption. For droughts, consider: crop yield losses + hydropower reduction + water-intensive industries."
], font_size=11, color=DARK_GRAY, spacing=Pt(6))

# Bottom: what the model does with it
add_shape_bg(slide, Inches(0.7), Inches(4.7), Inches(11.9), Inches(2.0), NEUTRAL_BG)
add_text_box(slide, Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.4),
             "What MINDSET does with the supply constraint",
             font_size=14, bold=True, color=DARK_BLUE)

add_text_box(slide, Inches(1.0), Inches(5.3), Inches(5.3), Inches(1.2),
             "1. Reads the constraint from the scenario file.\n"
             "2. Converts each entry to a sector-level output cap:\n"
             "       q_max = q_base \u00d7 (1 + Value)\n"
             "3. Applies the Ghosh inverse (G = (I\u2212B)\u207B\u00B9) to propagate\n"
             "   the supply shock forward through the IO system.",
             font_size=11, color=DARK_GRAY, font_name='Calibri')

add_text_box(slide, Inches(6.5), Inches(5.3), Inches(5.5), Inches(1.2),
             "4. Criticality weighting filters non-essential inputs:\n"
             "   critical (1.0) = full propagation, important (0.5) = 50%,\n"
             "   non-critical (0.0) = no propagation.\n"
             "5. Output reduction capped at \u221290% of baseline (floor).\n"
             "6. Result: dq_supply_constraint vector (output change by sector).",
             font_size=11, color=DARK_GRAY, font_name='Calibri')

add_source_line(slide,
    "Sources: Documentation/Capital damages.md (Hallegatte 2008, 2014a); "
    "SourceCode/InputOutput.py (lines 475\u2013560: calc_dq_supply_constraint; "
    "lines 140\u2013182: build_B_base with criticality)")

# ============================================================
# SLIDE 7: Inside MINDSET \u2014 What the Model Does
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "Inside MINDSET: Model Execution (RunMINDSET.py)",
             font_size=26, bold=True, color=WHITE)
add_footer(slide)

# Processing pipeline as numbered steps
proc_steps = [
    ("1. Load base data",
     "Loads the GLORIA MRIO database: intermediate transactions (IND_BASE), "
     "final demand (HH, Gov, GFCF), value added, employment, and sector criticality ratings. "
     "Entry point: exog_vars() in SourceCode/exog_vars.py.",
     VERY_LIGHT_BLUE),
    ("2. Read scenario file",
     "Parses the scenario Excel workbook. Reads supply constraints, final demand changes, "
     "investment, price changes, IO coefficient adjustments. Each sheet is optional. "
     "Entry point: scenario() in SourceCode/scenario.py.",
     LIGHT_GREEN_BG),
    ("3. Build Leontief and Ghosh inverses",
     "L = (I\u2212A)\u207B\u00B9 captures backward (demand-driven) linkages. "
     "G = (I\u2212B)\u207B\u00B9 captures forward (supply-driven) linkages, weighted by input criticality (Pichler et al. 2022). "
     "Entry point: IO.initialize() in SourceCode/InputOutput.py.",
     LIGHT_ORANGE_BG),
    ("4. Calculate all impact channels",
     "Supply constraint \u2192 Ghosh propagation. "
     "Investment/demand shocks \u2192 Leontief propagation. "
     "Price changes \u2192 trade and technology substitution. "
     "Each channel produces a separate dq (output change) vector.",
     NEUTRAL_BG),
    ("5. Iterate until convergence",
     "Output changes affect labor income \u2192 household consumption \u2192 demand \u2192 output. "
     "Tax revenue changes affect government spending. "
     "Loop repeats until labor compensation (<7%) and tax revenue (<1%) changes stabilize. "
     "Source: RunMINDSET.py lines 543\u2013688.",
     VERY_LIGHT_BLUE),
    ("6. Aggregate and save results",
     "Combines all 12 impact channels into totals. Decomposes output, employment, GDP, "
     "demand, household welfare, prices, revenue, and emissions by channel and sector. "
     "Saves to FullResults_[scenario].xlsx (8 sheets). Source: SourceCode/results.py.",
     LIGHT_GREEN_BG),
]

col1_x = Inches(0.7)
col2_x = Inches(6.8)
box_w = Inches(5.8)
box_h = Inches(1.6)
gap_y = Inches(0.15)

for i, (title, desc, bg) in enumerate(proc_steps):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = col1_x if col == 0 else col2_x
    y = Inches(1.3) + row * (box_h + gap_y)

    add_shape_bg(slide, x, y, box_w, box_h, bg)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), box_w - Inches(0.4), Inches(0.35),
                 title, font_size=13, bold=True, color=DARK_BLUE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), box_w - Inches(0.4), box_h - Inches(0.6),
                 desc, font_size=10.5, color=DARK_GRAY)

add_source_line(slide,
    "Sources: RunMINDSET.py (full pipeline); SourceCode/InputOutput.py "
    "(L_BASE, G_BASE, supply constraint); SourceCode/exog_vars.py (data loading)")

# ============================================================
# SLIDE 8: Model Outputs & Distributional Detail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "What You Get Back: Outputs and Distributional Detail",
             font_size=26, bold=True, color=WHITE)
add_footer(slide)

# Left: 8-sheet output
add_shape_bg(slide, Inches(0.7), Inches(1.3), Inches(5.0), Inches(4.0), LIGHT_GREEN_BG)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(4.5), Inches(0.4),
             "Results workbook (8 sheets)", font_size=15, bold=True, color=DARK_BLUE)

add_bullet_frame(slide, Inches(1.0), Inches(1.9), Inches(4.5), Inches(3.2), [
    "output: sectoral output changes (dq) by each of 12 impact channels",
    "employment: job changes by sector, decomposed by channel",
    "gdp: national GDP change by channel and region",
    "demand: final demand changes (HH, Gov, GFCF, etc.)",
    "household: consumption changes from price vs. income effects",
    "price: commodity price changes (pre/post trade substitution)",
    "revenue: tax revenue impacts by sector",
    "emissions: CO\u2082 and pollutant changes by sector",
], font_size=11, color=DARK_GRAY, spacing=Pt(4), bold_prefix=True)

# Right top: Employment detail
add_shape_bg(slide, Inches(5.9), Inches(1.3), Inches(6.7), Inches(1.8), VERY_LIGHT_BLUE)
add_text_box(slide, Inches(6.2), Inches(1.4), Inches(6.1), Inches(0.35),
             "Employment disaggregation", font_size=14, bold=True, color=DARK_BLUE)
add_bullet_frame(slide, Inches(6.2), Inches(1.8), Inches(6.1), Inches(1.1), [
    "2 skill levels (low-skill, high-skill) \u00d7 2 genders (female, male) = 4 labor categories",
    "Country-specific employment coefficients (jobs per $M output) for each of 120 sectors",
    "Wages tracked separately by skill level and gender",
], font_size=11, color=DARK_GRAY, spacing=Pt(3))

# Right bottom: Household detail
add_shape_bg(slide, Inches(5.9), Inches(3.3), Inches(6.7), Inches(2.0), LIGHT_ORANGE_BG)
add_text_box(slide, Inches(6.2), Inches(3.4), Inches(6.1), Inches(0.35),
             "Household consumption disaggregation", font_size=14, bold=True, color=DARK_BLUE)
add_bullet_frame(slide, Inches(6.2), Inches(3.8), Inches(6.1), Inches(1.3), [
    "3 income groups (low, middle, high) with differentiated cross-price elasticities (USDA)",
    "9 consumption categories: food, clothing, housing, furnishing, health, transport, recreation, education, other",
    "Price and income channels modeled separately \u2014 shows who bears the cost of commodity price increases",
    "Bolivia classified as Middle-income \u2014 middle-income elasticities apply across the country",
], font_size=11, color=DARK_GRAY, spacing=Pt(3))

# Bottom: what it does NOT capture
add_shape_bg(slide, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.2), LIGHT_RED_BG)
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.3),
             "What the distributional analysis does not capture",
             font_size=12, bold=True, color=ACCENT_RED)
add_bullet_frame(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.7), [
    "No microsimulation (no Gini, poverty headcount, or within-country inequality measures)",
    "No subnational resolution (cannot distinguish flood-prone lowlands from highlands)",
    "No informal economy (GLORIA based on formal national accounts; Bolivia's large informal sector is not captured)",
], font_size=10, color=DARK_GRAY, spacing=Pt(2))

add_source_line(slide,
    "Sources: SourceCode/results.py; README.md; "
    "Documentation/MINDSET modules/Impact channels/Household income and price.md; "
    "SourceCode/employment.py")

# ============================================================
# SLIDE 9: Bolivia Data Readiness
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "Bolivia: Data Availability & Readiness",
             font_size=28, bold=True, color=WHITE)
add_footer(slide)

items = [
    ("GLORIA MRIO (v57, 2019)",
     "Bolivia is 1 of 164 countries. Full IO table, bilateral trade, final demand, value added.",
     ACCENT_GREEN, "Available"),
    ("Employment coefficients",
     "Country-specific (Bolivia) employment-to-output ratios for all 120 sectors.",
     ACCENT_GREEN, "Available"),
    ("Labor data",
     "Disaggregated by skill (low/high) and gender (female/male) \u2014 volumes and wages.",
     ACCENT_GREEN, "Available"),
    ("Household income class",
     "Bolivia classified as Middle-income. USDA cross-price elasticities for middle-income group apply.",
     ACCENT_GREEN, "Available"),
    ("Energy own-price elasticities",
     "BOL-specific values in OwnPrices.xlsx.",
     ACCENT_GREEN, "Available"),
    ("Energy cross-price elasticities",
     "Uses Rest-of-World estimates (EnergyElas_ROW.xlsx), not Bolivia-specific.",
     ACCENT_ORANGE, "Generic"),
    ("Supply criticality matrix",
     "Global criticality ratings (Pichler et al. 2022) available, not Bolivia-customized.",
     ACCENT_ORANGE, "Generic"),
    ("Scenario files",
     "67 Bolivia-specific strategy scenario files already exist (Strategy_*_BOL.xlsx) as references.",
     ACCENT_GREEN, "Available"),
    ("Investment converter",
     "Defines how investment spending maps to sector demand. Available for all GLORIA countries.",
     ACCENT_GREEN, "Available"),
    ("Trade elasticities",
     "Default trade substitution elasticities (TradeElas_default.xlsx). Not Bolivia-specific.",
     ACCENT_ORANGE, "Generic"),
]

y = Inches(1.35)
for label, desc, color, status in items:
    badge_w = Inches(1.0)
    add_shape_bg(slide, Inches(0.7), y, badge_w, Inches(0.37), color)
    add_text_box(slide, Inches(0.7), y + Inches(0.03), badge_w, Inches(0.3),
                 status, font_size=10, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.85), y, Inches(3.0), Inches(0.37),
                 label, font_size=12, bold=True, color=DARK_GRAY)
    add_text_box(slide, Inches(4.85), y, Inches(7.8), Inches(0.37),
                 desc, font_size=11, color=DARK_GRAY)
    y += Inches(0.5)

add_source_line(slide,
    "Sources: GLORIA_template/Country_groupings/; GLORIA_template/Employment/Empl_coefficient.csv; "
    "GLORIA_db/v57/2019/; GLORIA_template/Elasticities/")

# ============================================================
# SLIDE 10: Task #2 \u2014 Energy Shocks
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "Task #2: Energy Shocks \u2014 What MINDSET Can and Cannot Do",
             font_size=28, bold=True, color=WHITE)
add_footer(slide)

# Can do
add_shape_bg(slide, Inches(0.7), Inches(1.4), Inches(5.8), Inches(5.3), LIGHT_GREEN_BG)
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(5.3), Inches(0.4),
             "What MINDSET can contribute", font_size=16, bold=True, color=ACCENT_GREEN)

add_bullet_frame(slide, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.4), [
    "Supply constraint on gas extraction (GLORIA sector 27) \u2014 propagates forward through all gas-dependent sectors",
    "Supply constraint on electricity (sector 63) \u2014 captures economy-wide impact of reduced power supply",
    "Interfuel substitution \u2014 energy cross-price elasticities model shifts between coal, oil, gas, electricity",
    "Price shock transmission \u2014 energy price increases propagate through all 120 sectors via IO linkages",
    "Macro outputs \u2014 GDP, employment, household welfare effects from externally defined energy shock magnitudes",
    "Energy sectors classified as essential \u2014 household/government consumption prioritized during rationing"
], font_size=12, color=DARK_GRAY, spacing=Pt(8))

# Cannot do
add_shape_bg(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.3), LIGHT_RED_BG)
add_text_box(slide, Inches(7.1), Inches(1.5), Inches(5.3), Inches(0.4),
             "Outside MINDSET\u2019s scope", font_size=16, bold=True, color=ACCENT_RED)

add_bullet_frame(slide, Inches(7.1), Inches(2.1), Inches(5.3), Inches(4.4), [
    "Resource depletion dynamics \u2014 no reserve curves, extraction cost functions, or exhaustion modeling",
    "Hydropower variability \u2014 no hydrological modeling, reservoir dynamics, or climate-water-energy nexus",
    "Energy system optimization \u2014 no generation capacity planning, least-cost expansion, or technology choice",
    "Low-carbon pathway investment \u2014 no learning curves, renewable cost declines, or structural transformation",
    "Technology-level electricity disaggregation \u2014 sector 63 aggregates all generation types (hydro, thermal, solar, wind)",
    "Recommended: pair with a dedicated energy system model (TIMES, OSeMOSYS, MESSAGE) for full Task #2 analysis"
], font_size=12, color=DARK_GRAY, spacing=Pt(8))

add_source_line(slide,
    "Sources: Documentation/MINDSET modules/Impact channels/Technology substitution.md; "
    "Documentation/Supply constraint.md; SourceCode/ener_elas.py")

# ============================================================
# SLIDE 11: Key Limitations & Caveats
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "Key Limitations & Caveats", font_size=28, bold=True, color=WHITE)
add_footer(slide)

# General
add_text_box(slide, Inches(0.7), Inches(1.4), Inches(5.8), Inches(0.4),
             "General Model Limitations", font_size=16, bold=True, color=DARK_BLUE)
add_bullet_frame(slide, Inches(0.9), Inches(1.9), Inches(5.4), Inches(4.5), [
    "No physical hazard modeling \u2014 all damage/shock magnitudes must be externally provided",
    "National aggregation only \u2014 no subnational spatial resolution",
    "Fixed IO coefficients \u2014 A and B matrices calibrated to 2019, not recalculated during simulation",
    "No inventories modeled \u2014 sectors cannot buffer shocks with stored inputs",
    "No firm heterogeneity within sectors",
    "Not a CGE \u2014 intentional for short-run analysis, limits long-run applicability",
    "Dynamic module is a prototype \u2014 variable source tables incomplete"
], font_size=12, color=DARK_GRAY, spacing=Pt(6))

# Bolivia-specific
add_text_box(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.4),
             "Bolivia-Specific Caveats", font_size=16, bold=True, color=DARK_BLUE)
add_bullet_frame(slide, Inches(7.0), Inches(1.9), Inches(5.4), Inches(4.5), [
    "Energy cross-price elasticities are generic (Rest-of-World), not Bolivia-specific",
    "Household elasticities use USDA data \u2014 quality for Bolivia not independently verified",
    "Large informal economy not captured in MRIO-based national accounts",
    "Agricultural vulnerability to ENSO requires external mapping (which crops are drought/flood sensitive)",
    "Employment coefficients from 2019 \u2014 may not reflect recent structural shifts",
    "Country-level results \u2014 cannot distinguish between Altiplano, valleys, and lowlands"
], font_size=12, color=DARK_GRAY, spacing=Pt(6))

add_source_line(slide,
    "Sources: Documentation/Supply constraint.md; Documentation/Capital damages.md; "
    "GLORIA_template/Country_groupings/")

# ============================================================
# SLIDE 12: Summary Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
             "Summary: Capability vs. Need", font_size=28, bold=True, color=WHITE)
add_footer(slide)

headers = ["Analysis Need", "MINDSET Capability", "Gap"]
col_widths = [Inches(4.2), Inches(4.2), Inches(3.4)]
col_starts = [Inches(0.7), Inches(4.9), Inches(9.1)]

y_start = Inches(1.4)
row_h = Inches(0.45)

for j, (header, x, w) in enumerate(zip(headers, col_starts, col_widths)):
    add_shape_bg(slide, x, y_start, w, row_h, DARK_BLUE)
    add_text_box(slide, x + Inches(0.1), y_start + Inches(0.05), w - Inches(0.2), row_h,
                 header, font_size=12, bold=True, color=WHITE)

rows = [
    ("ENSO physical hazard characterization",
     "None \u2014 requires external climate/hydrology model", "Full gap"),
    ("Sector-level damage estimates",
     "None \u2014 requires external damage assessment", "Full gap"),
    ("Macroeconomic impact propagation",
     "Full support (GDP, output, prices, all 12 channels)", "None"),
    ("Employment by sector / skill / gender",
     "Full support (4 labor categories \u00d7 120 sectors)", "None"),
    ("Household distributional effects",
     "Partial \u2014 3 income groups, elasticity-based", "No inequality / poverty"),
    ("Subnational spatial impacts",
     "None", "Full gap"),
    ("Short / medium-term time horizon",
     "Supported (static + prototype dynamic)", "Dynamic is prototype"),
    ("Gas reserve depletion dynamics",
     "None \u2014 needs resource economics model", "Full gap"),
    ("Hydropower variability",
     "None \u2014 needs energy/hydrology model", "Full gap"),
    ("Energy system optimization / pathways",
     "None \u2014 needs TIMES / OSeMOSYS", "Full gap"),
    ("Energy shock macro propagation",
     "Full support (supply constraint + price channels)", "None"),
]

for i, (need, capability, gap) in enumerate(rows):
    y = y_start + (i + 1) * row_h
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, (text, x, w) in enumerate(zip([need, capability, gap], col_starts, col_widths)):
        add_shape_bg(slide, x, y, w, row_h, bg)
        clr = (ACCENT_GREEN if "Full support" in text
               else (ACCENT_RED if "Full gap" in text else DARK_GRAY))
        add_text_box(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), row_h,
                     text, font_size=10,
                     bold=("Full support" in text or "Full gap" in text),
                     color=clr)

# ============================================================
# Save
# ============================================================
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "MINDSET_Bolivia_Assessment.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
